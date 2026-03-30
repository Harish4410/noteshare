"""
utils/ai_utils.py - AI features using Google Gemini + text extraction
"""
import os, json, re
import urllib.request

GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', '').strip()
GEMINI_URL     = "https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent"

# ── Gemini API call ───────────────────────────────────────────
def _call_gemini(prompt: str) -> str:
    if not GEMINI_API_KEY:
        return ''
    try:
        payload = json.dumps({
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {"temperature": 0.4, "maxOutputTokens": 1024}
        }).encode()
        req = urllib.request.Request(
            f"{GEMINI_URL}?key={GEMINI_API_KEY}",
            data=payload,
            headers={"Content-Type": "application/json"},
            method="POST"
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            data = json.loads(resp.read())
        return data['candidates'][0]['content']['parts'][0]['text'].strip()
    except Exception as e:
        print(f"[AI] Gemini error: {e}")
        return ''

def _use_gemini():
    return bool(GEMINI_API_KEY)

# ── Text Extraction ───────────────────────────────────────────
def extract_text(file_path: str) -> str:
    """Extract text from PDF, DOCX, DOC, TXT, PPT, PPTX files."""
    if not os.path.exists(file_path):
        print(f"[AI] File not found: {file_path}")
        return ''
    ext = file_path.rsplit('.', 1)[-1].lower() if '.' in file_path else ''
    try:
        if ext == 'pdf':
            return _extract_pdf(file_path)
        elif ext in ('docx', 'doc'):
            return _extract_docx(file_path)
        elif ext in ('pptx', 'ppt'):
            return _extract_pptx(file_path)
        elif ext == 'txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext in ('png', 'jpg', 'jpeg'):
            return ''  # Images can't be text-extracted without OCR
        else:
            # Try reading as plain text anyway
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except:
                return ''
    except Exception as e:
        print(f'[AI] Extract error ({ext}): {e}')
        return ''

def _extract_pdf(path):
    # Try pypdf first
    try:
        import pypdf
        text = []
        with open(path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                t = page.extract_text()
                if t: text.append(t)
        result = '\n'.join(text)
        if result.strip():
            return result
    except Exception as e:
        print(f'[AI] pypdf error: {e}')

    # Fallback: PyPDF2
    try:
        import PyPDF2
        text = []
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                t = page.extract_text()
                if t: text.append(t)
        return '\n'.join(text)
    except Exception as e:
        print(f'[AI] PyPDF2 error: {e}')
        return ''

def _extract_docx(path):
    # Try python-docx
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return '\n'.join(paragraphs)
    except ImportError:
        print('[AI] python-docx not installed. Run: pip install python-docx')
    except Exception as e:
        print(f'[AI] DOCX error: {e}')

    # Fallback: try reading as zip and extracting XML
    try:
        import zipfile, xml.etree.ElementTree as ET
        with zipfile.ZipFile(path) as z:
            with z.open('word/document.xml') as f:
                tree  = ET.parse(f)
                texts = tree.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t')
                return ' '.join(t.text for t in texts if t.text)
    except Exception as e:
        print(f'[AI] DOCX fallback error: {e}')
        return ''

def _extract_pptx(path):
    try:
        from pptx import Presentation
        prs   = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
        return '\n'.join(texts)
    except Exception as e:
        print(f'[AI] PPTX error: {e}')
        return ''

# ── Summary ───────────────────────────────────────────────────
def generate_summary(text: str) -> str:
    text = text[:4000]
    if _use_gemini():
        result = _call_gemini(
            f"Summarize these academic notes in exactly 5 clear bullet points. "
            f"Each bullet must start with •\n\n{text}"
        )
        if result: return result

    sentences = re.split(r'(?<=[.!?])\s+', text)
    return '\n'.join(f'• {s.strip()}' for s in sentences[:5] if len(s.strip()) > 20)

# ── Flashcards ────────────────────────────────────────────────
def generate_flashcards(text: str) -> list:
    text = text[:3500]
    if _use_gemini():
        result = _call_gemini(
            f"Create 8 flashcards from these notes. "
            f"Return ONLY a valid JSON array, no extra text. "
            f'Format: [{{"question":"...","answer":"..."}}]\n\n{text}'
        )
        if result:
            try:
                clean = re.sub(r'```json|```', '', result).strip()
                cards = json.loads(clean)
                if isinstance(cards, list) and cards:
                    return cards
            except: pass

    # Fallback
    cards = []
    for m in re.finditer(
        r'([A-Z][^.?!]{10,80})\s+(?:is|are|means?|refers? to)\s+([^.?!]{10,150})', text
    ):
        cards.append({'question': f'What {m.group(1).lower()}?',
                      'answer':   m.group(2).strip()})
        if len(cards) >= 8: break
    if not cards:
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text)
                     if len(s.strip()) > 30][:8]
        for s in sentences:
            words = s.split(); mid = max(1, len(words)//2)
            cards.append({'question': ' '.join(words[:mid]) + ' ___?',
                          'answer':   ' '.join(words[mid:])})
    return cards or [{'question': 'Review this note', 'answer': text[:100]}]

# ── Quiz ──────────────────────────────────────────────────────
def generate_quiz(text: str) -> list:
    text = text[:3500]
    if _use_gemini():
        result = _call_gemini(
            f"Create 5 multiple choice questions from these notes. "
            f"Return ONLY a valid JSON array, no extra text. "
            f'Format: [{{"question":"...","options":["A","B","C","D"],"answer":"correct text"}}]\n\n{text}'
        )
        if result:
            try:
                clean = re.sub(r'```json|```', '', result).strip()
                quiz  = json.loads(clean)
                if isinstance(quiz, list) and quiz:
                    return quiz
            except: pass

    return [{'question': 'What is the main topic of these notes?',
             'options': ['Option A', 'Option B', 'Option C', 'Option D'],
             'answer': 'Option A'}]

# ── Score / Evaluate ──────────────────────────────────────────
def evaluate_note(text: str) -> float:
    if _use_gemini():
        result = _call_gemini(
            f"Rate the quality of these academic notes from 0 to 10. "
            f"Consider clarity, completeness, structure, usefulness. "
            f"Reply with ONLY a single number like 7.5\n\n{text[:3000]}"
        )
        if result:
            try:
                m = re.search(r'\d+(\.\d+)?', result)
                if m: return round(min(10.0, max(0.0, float(m.group()))), 1)
            except: pass

    wc   = len(text.split())
    hdrs = len(re.findall(r'\n[A-Z][A-Z\s]{3,}', text))
    return min(10.0, max(1.0, round((wc / 200) + (hdrs * 0.5), 1)))
